from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os
from datetime import datetime

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.slide_width = 13.333
        self.slide_height = 7.5
        self.margin = 1.0
        self.grid_spacing = 0.25

    def _create_title_slide(self, title="Patient Treatment Report"):
        """Create the title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = "Generated on " + datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    def _create_image_grid_slide(self, images, category_title):
        """Create a slide with a left-aligned 2x2 grid of images."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Set fixed dimensions for the grid
        grid_width = 6.0
        grid_height = 4.0
        
        # Test with extreme left position
        grid_left = 2.0  # Fixed 2 inches from left
        grid_top = 2.0   # Fixed 2 inches from top

        # Always add header (category title)
        title_box = slide.shapes.add_textbox(
            Inches(0.5),  # Left margin
            Inches(0.5),  # Top margin
            Inches(12.0),  # Width
            Inches(0.8)   # Height
        )
        title_frame = title_box.text_frame
        title_frame.text = category_title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        # Add placeholder text only for non-unknown categories
        if category_title.lower() != "unknown":
            text_box = slide.shapes.add_textbox(
                Inches(0.5),  # Left margin
                Inches(1.3),  # Below title
                Inches(12.0), # Width
                Inches(0.5)   # Height
            )
            text_frame = text_box.text_frame
            text_frame.text = "The following visuals provide an overview of treatment outcomes"
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Calculate individual image dimensions
        img_width = (grid_width - self.grid_spacing) / 2
        img_height = (grid_height - self.grid_spacing) / 2

        # Debug print statements
        print("\nDebug Layout Information:")
        print(f"Slide dimensions: {self.slide_width} x {self.slide_height} inches")
        print(f"Grid dimensions: {grid_width} x {grid_height} inches")
        print(f"Grid position: left={grid_left}, top={grid_top} inches")
        print(f"Image dimensions: {img_width} x {img_height} inches")
        print(f"Grid spacing: {self.grid_spacing} inches")

        # Add images to the grid
        for idx, image_data in enumerate(images):
            if idx >= 4:  # Maximum 4 images per slide
                break
                
            row = idx // 2
            col = idx % 2
            
            # Simplified positioning for testing
            left = grid_left + (col * (img_width + self.grid_spacing))
            top = grid_top + (row * (img_height + self.grid_spacing))
            
            # Debug print statements for each image
            print(f"\nImage {idx + 1} position:")
            print(f"Row: {row}, Column: {col}")
            print(f"Position: left={left}, top={top} inches")
            
            # Add the image with explicit size
            slide.shapes.add_picture(
                image_data['path'],
                Inches(left),
                Inches(top),
                Inches(img_width),
                Inches(img_height)
            )

    def generate_presentation(self, categorized_images, output_path):
        """Generate a PowerPoint presentation from categorized images."""
        # Create title slide
        self._create_title_slide()
        
        # Define the desired order of categories
        category_order = [
            'intra_oral',
            'front_with_teeth',
            'front_no_teeth',
            'side_view',
            'unknown'  # Keep unknown at the end
        ]
        
        # Create slides for each category in the specified order
        for category in category_order:
            images = categorized_images.get(category, [])
            if not images:
                continue
                
            # Sort images by creation date and time
            sorted_images = sorted(
                images,
                key=lambda x: x['metadata'].get('creation_date', datetime.min)
            )
            
            # Create slides with 4 images each
            for i in range(0, len(sorted_images), 4):
                slide_images = sorted_images[i:i+4]
                self._create_image_grid_slide(
                    slide_images,
                    images[0]['category_label'] if images else category.replace('_', ' ').title()
                )
        
        # Save the presentation
        self.prs.save(output_path)
        return output_path

    def _resize_image(self, image_path, max_size=(1920, 1080)):
        """Resize an image while maintaining aspect ratio."""
        with Image.open(image_path) as img:
            img.thumbnail(max_size, Image.Resampling.LANCZOS)
            return img 